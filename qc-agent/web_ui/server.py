import sys
import os
# Add project root to path first
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

# Set NO_PROXY immediately
os.environ['NO_PROXY'] = 'localhost,127.0.0.1'
os.environ['no_proxy'] = 'localhost,127.0.0.1'

import urllib.request
import subprocess
import tempfile
import platform
print("DEBUG: Proxy environment is using:", urllib.request.getproxies())

from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename
import requests
from bs4 import BeautifulSoup
import pypdf
import docx
import openpyxl
import pptx




# --- Pure-Python text splitter (thay thế langchain_text_splitters) ---
def _split_text_recursive(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[str]:
    """Chia text thành các chunks có kích thước chunk_size với overlap chunk_overlap.
    Ưu tiên tách tại: \\n\\n -> \\n -> ' ' -> ký tự.
    """
    separators = ["\n\n", "\n", " ", ""]

    def _split(t: str, seps: list[str]) -> list[str]:
        if not seps:
            return [t] if t else []
        sep = seps[0]
        rest = seps[1:]
        parts = t.split(sep) if sep else list(t)
        result = []
        for part in parts:
            if len(part) > chunk_size:
                result.extend(_split(part, rest))
            elif part:
                result.append(part)
        return result

    raw_pieces = _split(text, separators)

    chunks: list[str] = []
    current = ""
    for piece in raw_pieces:
        candidate = (current + "\n" + piece).strip() if current else piece
        if len(candidate) <= chunk_size:
            current = candidate
        else:
            if current:
                chunks.append(current)
            # overlap: lấy phần cuối của current làm đầu chunk mới
            overlap_text = current[-chunk_overlap:] if chunk_overlap and current else ""
            current = (overlap_text + "\n" + piece).strip() if overlap_text else piece
    if current:
        chunks.append(current)
    return chunks


try:
    import config
    from memory.longterm_memory import LongTermMemory
    from src.utils import jira_client
    from src.utils.sandbox_docker import start_docker_sandbox, SandboxConfig
except ImportError as e:
    print(f"❌ Lỗi nghiêm trọng: Không thể import các module cần thiết của dự án. {e}")
    sys.exit(1)

# --- Khởi tạo Flask App ---
UI_DIRECTORY = os.path.dirname(os.path.abspath(__file__))

app = Flask(__name__)
CORS(app)
app.config['UPLOAD_FOLDER'] = './uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# --- Khởi tạo Bộ Nhớ Dài Hạn ---
long_term_memory = None
try:
    print("🧠 Đang khởi tạo Bộ nhớ dài hạn cho máy chủ...")
    long_term_memory = LongTermMemory()
    print("✅ Khởi tạo Bộ nhớ dài hạn thành công.")
except Exception as e:
    print(f"❌ Không thể khởi tạo LongTermMemory: {e}")

# --- BIẾN TOÀN CỤC CHO TRẠNG THÁI HỆ THỐNG ---
# Lưu trữ trạng thái của 10 lần thực thi gần nhất (PASS/FAIL)
# Cấu trúc: [{"id": "...", "status": "PASS/FAIL", "timestamp": "..."}]
execution_history = []

# --- Các Hàm Hỗ Trợ Trích Xuất Văn Bản ---
def extract_text_from_url(url: str) -> str:
    # ... (code không thay đổi)
    try:
        response = requests.get(url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        text = soup.get_text(separator='\n', strip=True)
        return text
    except Exception as e:
        raise ValueError(f"Không thể lấy nội dung từ URL: {e}")

def extract_text_from_pdf(filepath: str) -> str:
    with open(filepath, 'rb') as f:
        reader = pypdf.PdfReader(f)
        return "\n".join(page.extract_text() for page in reader.pages)

def extract_text_from_docx(filepath: str) -> str:
    doc = docx.Document(filepath)
    return "\n".join(para.text for para in doc.paragraphs)

def extract_text_from_xlsx(filepath: str) -> str:
    workbook = openpyxl.load_workbook(filepath)
    text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value:
                    text.append(str(cell.value))
    return "\n".join(text)

def extract_text_from_pptx(filepath: str) -> str:
    prs = pptx.Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)
    
def extract_text_from_txt(filepath: str) -> str:
    with open(filepath, 'r', encoding='utf-8') as f:
        return f.read()

# --- Các API Endpoints ---

@app.route('/')
def serve_index():
    return send_from_directory(UI_DIRECTORY, 'index.html')

@app.route('/info', methods=['GET'])
def get_system_info():
    """Endpoint mới: Cung cấp thông tin chi tiết về hệ thống."""
    if not long_term_memory:
        return jsonify({"error": "[STATUS]: LONG_TERM_MEMORY_NOT_INITIALIZED"}), 500
    
    embedding_model = (
        config.EMBEDDING_MODEL_NAME_OPENAI 
        if config.EMBEDDING_API_PROVIDER == "openai" 
        else config.EMBEDDING_MODEL_NAME_LOCAL
    )
    
    info = {
        "provider": config.VECTOR_DB_PROVIDER,
        "embedding_provider": config.EMBEDDING_API_PROVIDER,
        "embedding_model": embedding_model,
        "item_count": long_term_memory.count() if long_term_memory else 0
    }
    return jsonify(info)

@app.route('/models', methods=['GET'])
def get_available_models():
    """Trả về danh sách các model khả dụng cho UI."""
    return jsonify({
        "models": getattr(config, 'AVAILABLE_MODELS', ["Qwen3-32B"]),
        "default_planner": getattr(config, 'FAST_PLANNER_MODEL', "Qwen3-32B"),
        "default_tester": getattr(config, 'SENIOR_TESTER_MODEL', "Qwen3-32B")
    })

@app.route('/retrieve', methods=['POST'])
def retrieve_chunks():
    """Endpoint mới: Thực hiện truy xuất từ vector DB."""
    if not long_term_memory:
        return jsonify({"error": "[STATUS]: LTM_DISCONNECTED_EXCEPTION"}), 500

    data = request.get_json()
    query = data.get('query')
    top_k = data.get('top_k', 3)

    if not query:
        return jsonify({"error": "Câu truy vấn không được để trống."}), 400

    try:
        results = long_term_memory.retrieve(query=query, top_k_broad=int(top_k), top_k_final=int(top_k))
        return jsonify({"results": results})
    except Exception as e:
        return jsonify({"error": f"Lỗi khi truy xuất: {e}"}), 500

@app.route('/clear-collection', methods=['POST'])
def clear_collection():
    """Xóa toàn bộ dữ liệu trong vector DB collection."""
    if not long_term_memory:
        return jsonify({"error": "Bộ nhớ dài hạn chưa được khởi tạo."}), 500
    try:
        long_term_memory.db.clear_collection()
        return jsonify({"message": "Collection đã được xóa và tạo lại thành công."}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

def process_and_add_text(text: str, source: str):
    # ... (code không thay đổi)
    if not long_term_memory:
        raise ConnectionError("Máy chủ chưa kết nối được với Vector DB.")
    if not text.strip():
        raise ValueError("Không tìm thấy nội dung văn bản để thêm.")
    
    chunks = _split_text_recursive(text, chunk_size=1000, chunk_overlap=200)

    print(f"📄 Văn bản gốc được chia thành {len(chunks)} chunks.")

    for i, chunk_text in enumerate(chunks):
        chunk_with_metadata = f"[Nguồn: {source}, Phần: {i+1}/{len(chunks)}]\n\n{chunk_text}"
        long_term_memory.add_memory(chunk_with_metadata)

    total_count = long_term_memory.db.count()
    status_msg = f"Thêm thành công {len(chunks)} phần từ '{source}'. Tổng số ký ức: {total_count}"
    return status_msg, len(chunks)

@app.route('/add-text', methods=['POST'])
def add_text():
    # ... (code không thay đổi)
    try:
        data = request.get_json()
        text = data['text']
        message, count = process_and_add_text(text, "Văn bản dán trực tiếp")
        return jsonify({"message": message, "chunk_count": count}), 201
    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route('/add-url', methods=['POST'])
def add_url():
    # ... (code không thay đổi)
    try:
        data = request.get_json()
        url = data['url']
        text = extract_text_from_url(url)
        message, count = process_and_add_text(text, f"URL: {url}")
        return jsonify({"message": message, "chunk_count": count}), 201
    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route('/add-file', methods=['POST'])
def add_file():
    # ... (code không thay đổi)
    if 'file' not in request.files:
        return jsonify({"error": "Không tìm thấy tệp trong yêu cầu."}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Chưa chọn tệp nào."}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)

    try:
        ext = os.path.splitext(filename)[1].lower()
        text = ""
        if ext == '.pdf':
            text = extract_text_from_pdf(filepath)
        elif ext == '.docx':
            text = extract_text_from_docx(filepath)
        elif ext == '.xlsx':
            text = extract_text_from_xlsx(filepath)
        elif ext == '.pptx':
            text = extract_text_from_pptx(filepath)
        elif ext in ['.txt', '.md']:
            text = extract_text_from_txt(filepath)
        else:
            raise ValueError(f"Định dạng tệp '{ext}' không được hỗ trợ.")
        
        message, count = process_and_add_text(text, f"Tệp: {filename}")
        return jsonify({"message": message, "chunk_count": count}), 201
    except Exception as e:
        return jsonify({"error": str(e)}), 400
    finally:
        if os.path.exists(filepath):
            os.remove(filepath)


@app.route('/execution-history', methods=['GET'])
def get_execution_history():
    """Trả về 10 lần thực thi gần nhất để hiển thị AI Agent Status Bar."""
    return jsonify(execution_history[-10:])

import asyncio
from fast_track_endpoint import run_fast_track_logic

@app.route('/fast-track', methods=['POST'])
def fast_track_test():
    """
    Endpoint Fast Track Testing (Simplified Wrapper).
    Delegates logic to fast_track_endpoint.run_fast_track_logic.
    """
    try:
        data = request.get_json()
        
        # 1. Execute consolidated logic (Use asyncio.run for sync compatibility)
        result = asyncio.run(run_fast_track_logic(data))
        
        if "error" in result:
            return jsonify(result), 500

        # 2. Update execution history for UI status bar
        from datetime import datetime
        overall_status = result.get('overall_status', 'FAIL')
        execution_history.append({
            "id": f"FT-{int(datetime.now().timestamp())}",
            "status": overall_status,
            "timestamp": datetime.now().isoformat()
        })
        if len(execution_history) > 20:
            execution_history.pop(0)

        return jsonify(result), 200
        
    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        print(error_msg)
        with open("tmp/server_error.log", "w", encoding="utf-8") as f:
            f.write(error_msg)
        return jsonify({"error": f"[STATUS]: TRACE_EXCEPTION_{type(e).__name__.upper()}", "details": str(e)}), 500

@app.route('/push-to-jira', methods=['POST'])
def push_to_jira():
    """Duyệt và đẩy nội dung lên Jira."""
    data = request.json
    # 1. Handle Bulk Push if test_cases list is provided
    test_cases = data.get('test_cases', [])
    if test_cases:
        results = []
        for tc_data in test_cases:
            res = _push_single_case(tc_data)
            results.append(res)
        return jsonify({"results": results}), 200
    
    # 2. Handle Single Push (Legacy/Manual)
    res = _push_single_case(data)
    if "error" in res:
        return jsonify(res), 500
    return jsonify(res), 200

def _push_single_case(data):
    task_name = data.get('target', data.get('description', 'New Task'))
    worker_evidence = data.get('evidence', '')
    status = data.get('status', 'PASS')
    cycle_id = data.get('cycle_id')
    case_id = data.get('jira_id', data.get('jira_test_case_id'))
    project_key = data.get('project_key', 'NCPP')

    # 1. Nếu chưa có mã Case ID (NCPP-T...), tiến hành tạo mới
    if not case_id:
        print(f"[DEBUG] 🆕 Duyệt tạo mới Test Case cho: {task_name}")
        case_id = jira_client.create_zephyr_test_case(
            name=f"[{data.get('tester', 'Tester')}] {task_name}",
            objective=worker_evidence,
            project_key=project_key,
            folder=data.get('folder', '')
        )
        if not case_id:
            return {"error": f"Không thể tạo Test Case cho {task_name}"}
    
    # 2. Nếu có Cycle ID, tiến hành đồng bộ kết quả thực thi (Test Execution)
    if case_id and cycle_id:
        automation_code = data.get('automation_code', '')
        execution_log = f"Status: {status}\n\nEvidence Summary:\n{worker_evidence[:1000]}"
        if automation_code:
            execution_log += f"\n\n[Automation Code]\n{automation_code}"
            
        print(f"[DEBUG] 🔄 Đang đẩy kết quả thực thi lên Jira cho {case_id} trong Cycle {cycle_id}...")
        success = jira_client.sync_and_execute_test_case(
            cycle_id=cycle_id,
            test_case_id=case_id,
            status=status,
            log=execution_log
        )
        if not success:
            print("⚠️ Failed to create execution record for {case_id}")
            return {"case_id": case_id, "status": "FAIL", "message": "Execution sync failed"}
        else:
            print(f"✅ Successfully synced {status} result to Jira for {case_id}.")
            return {"case_id": case_id, "status": "PASS"}
    
    return {"case_id": case_id, "status": "UNKNOWN"}

    return jsonify({
        "status": "success",
        "jira_test_case_id": case_id,
        "message": f"Đã đẩy thành công lên Jira với mã {case_id}"
    })

@app.route('/extract-cycle', methods=['POST'])
def extract_cycle():
    """Trích xuất danh sách Test Case từ Cycle."""
    data = request.json
    cycle_id = data.get('cycle_id')
    
    if not cycle_id:
        return jsonify({"error": "Vui lòng nhập Cycle ID"}), 400
        
    try:
        print(f"[DEBUG] 🔍 Đang trích xuất toàn bộ Test Case từ Cycle: {cycle_id}...")
        cycle_data = jira_client.get_test_cycle_info(cycle_id)
        
        # Trích xuất danh sách gọn nhẹ
        items = cycle_data.get('testRunItems', {}).get('testRunItems', [])
        cases = []
        for item in items:
            # item is a dict, $lastTestResult is a key
            last_res = item.get('$lastTestResult', {})
            tc = last_res.get('testCase', {}) if last_res else {}
            
            cases.append({
                "key": tc.get('key', 'Unknown'),
                "name": tc.get('name', 'N/A'),
                "id": tc.get('id'),
                "status": last_res.get('testResult', {}).get('name', 'N/A') if last_res else 'N/A'
            })
            
        return jsonify({
            "cycle_id": cycle_id,
            "project_name": cycle_data.get('projectId'),
            "total_cases": len(cases),
            "cases": cases
        })
    except Exception as e:
        print(f"❌ Lỗi trích xuất Cycle: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/sync-cycle-to-vector', methods=['POST'])
def sync_cycle_to_vector():
    """Trích xuất chi tiết toàn bộ Cycle, ghi file và đưa vào Vector DB."""
    data = request.json
    cycle_id = data.get('cycle_id')
    
    if not cycle_id or not long_term_memory:
        return jsonify({"error": "Thiếu Cycle ID hoặc Vector DB chưa sẵn sàng"}), 400
        
    try:
        print(f"[DEBUG] 🚀 Bắt đầu đồng bộ Cycle {cycle_id} vào Vector DB...")
        cycle_data = jira_client.get_test_cycle_info(cycle_id)
        items = cycle_data.get('testRunItems', {}).get('testRunItems', [])
        
        full_content = []
        count = 0
        
        for item in items:
            last_res = item.get('$lastTestResult', {})
            tc = last_res.get('testCase', {}) if last_res else {}
            case_id = tc.get('key')
            if not case_id: continue
            
            # 1. Fetch chi tiết từng Case (vẫn phải loop vì API lấy từng case)
            detail = jira_client.get_test_case_info(case_id)
            print(f"[DEBUG] Fetched test case info for {case_id}")
            
            steps_text = ""
            script = detail.get('testScript', {}).get('stepByStepScript', {})
            if script:
                for s in script.get('steps', []):
                    steps_text += f"\n- Step {s.get('index', 0)+1}: {s.get('description')} => {s.get('expectedResult')}"
            
            case_block = f"TEST CASE ID: {case_id}\nNAME: {tc.get('name')}\nOBJECTIVE: {detail.get('objective', '')}\nSTEPS: {steps_text}\n---\n"
            full_content.append(case_block)
            count += 1

        # 2. Nạp HÀNG LOẠT vào Vector DB sau khi đã fetch xong
        if full_content:
            print(f"[DEBUG] 📥 Đang nạp hàng loạt {len(full_content)} kịch bản vào Vector DB...")
            long_term_memory.add_memories(full_content)

        # Ghi ra file vật lý để lưu trữ
        filename = f"{cycle_id}_full_dump.txt"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write("\n".join(full_content))
            
        return jsonify({
            "status": "success",
            "message": f"Đã đồng bộ {count} kịch bản vào Vector DB và lưu file {filename}",
            "file": filename
        })
    except Exception as e:
        print(f"❌ Lỗi đồng bộ: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/list-uploads', methods=['GET'])
def list_uploads():
    """Liệt kê các file log trong bộ nhớ uploads."""
    files = os.listdir(app.config['UPLOAD_FOLDER'])
    return jsonify({"files": sorted(files, reverse=True)})

@app.route('/view-upload/<filename>', methods=['GET'])
def view_upload(filename):
    """Xem nội dung chi tiết của một file log."""
    try:
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()
        return jsonify({"content": content})
    except Exception as e:
        return jsonify({"error": str(e)}), 404

@app.route('/execute-code', methods=['POST'])
def execute_code():
    """Thực thi mã Python/Playwright tại chỗ và trả về kết quả console."""
    data = request.json
    code = data.get('code', '')
    
    if not code:
        return jsonify({"error": "Không có mã để thực thi"}), 400
        
    try:
        # Tạo file tạm để thực thi
        # Lưu ý: sandbox mount /tmp nên ta nên dùng /tmp trên host để container thấy được
        tmp_dir = "/tmp" if platform.system() != "Windows" else tempfile.gettempdir()
        with tempfile.NamedTemporaryFile(suffix=".py", delete=False, dir=tmp_dir, mode='w', encoding='utf-8') as f:
            f.write(code)
            temp_path = f.name
            
        print(f"[DEBUG] ⚡ Đang chuẩn bị thực thi mã tại: {temp_path}")
        
        # Kiểm tra xem có sử dụng sandbox không
        use_sandbox = os.environ.get("SANDBOX") == "1"
        
        if use_sandbox:
            print("[DEBUG] 🐳 Đang khởi chạy trong Docker Sandbox...")
            # Chuyển path host sang path container
            from src.utils.sandbox_docker import get_container_path
            container_temp_path = get_container_path(temp_path)
            
            image = os.environ.get("SANDBOX_IMAGE", "qc-agent-framework:latest")
            sandbox_cmd = os.environ.get("SANDBOX_COMMAND", "docker")
            
            config_obj = SandboxConfig(command=sandbox_cmd, image=image)
            cli_cmd = f"python3 {container_temp_path}"
            
            result = asyncio.run(start_docker_sandbox(cli_cmd, config_obj, capture_output=True))
            
            # Xóa file tạm sau khi chạy
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
            return jsonify({
                "stdout": result["stdout"],
                "stderr": result["stderr"],
                "exit_code": result["exit_code"],
                "sandbox": True
            })
        else:
            # Chạy code local với subprocess như cũ
            print("[DEBUG] 🏠 Đang chạy trực tiếp trên host...")
            process = subprocess.run(
                [sys.executable, temp_path],
                capture_output=True,
                text=True,
                timeout=60,
                env=os.environ.copy()
            )
            
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
            return jsonify({
                "stdout": process.stdout,
                "stderr": process.stderr,
                "exit_code": process.returncode,
                "sandbox": False
            })
        
    except subprocess.TimeoutExpired:
        if 'temp_path' in locals() and os.path.exists(temp_path): os.remove(temp_path)
        return jsonify({"error": "Thực thi quá thời gian (Timeout 60s)"}), 500
    except Exception as e:
        if 'temp_path' in locals() and os.path.exists(temp_path): os.remove(temp_path)
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    print("🚀 Bắt đầu máy chủ Flask cho Giao diện VectorDB...")
    app.run(host='0.0.0.0', port=5001, debug=True)
